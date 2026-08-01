"""
Document -> list of chunks.

Design choices (see README for the full rationale):
- We split on paragraph/heading boundaries first, THEN pack paragraphs into
  token-bounded windows with overlap. Splitting on fixed character counts
  cuts sentences in half and hurts retrieval quality; splitting purely on
  paragraphs without a size cap gives wildly uneven chunk sizes. Packing
  paragraphs into a token budget gets both semantic coherence and
  consistent chunk size.
- Overlap is applied in tokens between consecutive chunks so a fact that
  sits right on a paragraph boundary is still fully present in at least
  one chunk.
- We keep a per-page (or per-slide/per-sheet) text map for every format so
  every chunk can be tagged with a location reference, which is what the
  "source" citations in the UI point back to. For formats with no native
  concept of a page (docx, plain text), page_number is None and citations
  fall back to section title only.
- Extraction is format-dispatched by file extension in extract_pages(). Each
  branch is intentionally isolated so adding a new format later means adding
  one branch, not touching chunking/packing logic at all.
"""
import io
import re
from dataclasses import dataclass

import tiktoken
from pypdf import PdfReader

_ENCODER = tiktoken.get_encoding("cl100k_base")

# Extensions the upload endpoint accepts. Kept here (not in the router) so the
# list of "what we can actually extract text from" and "what we advertise as
# accepted" can never drift apart.
SUPPORTED_EXTENSIONS = {
    "pdf", "txt", "md", "docx", "pptx", "csv", "xlsx", "xls",
    "html", "htm", "png", "jpg", "jpeg", "webp",
}


def count_tokens(text: str) -> int:
    return len(_ENCODER.encode(text))


@dataclass
class RawPage:
    page_number: int | None  # None when the format has no native page concept
    text: str


@dataclass
class Chunk:
    content: str
    token_count: int
    page_number: int | None
    section_title: str | None


def _file_ext(filename: str) -> str:
    return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""


def _ocr_image_bytes(image_bytes: bytes) -> str:
    """OCR a single raster image. Isolated behind one function so the
    tesseract/Pillow dependency is only imported when actually needed."""
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(image_bytes))
    return pytesseract.image_to_string(img) or ""


def _extract_pdf(file_bytes: bytes) -> list[RawPage]:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages: list[RawPage] = []
    ocr_doc = None  # lazily opened only if a page needs OCR fallback

    for i, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()

        if not text:
            # Likely a scanned/image-only page — fall back to OCR by
            # rendering the page to an image with PyMuPDF and running
            # tesseract over it. Only pay this cost for pages that actually
            # need it.
            try:
                import fitz  # PyMuPDF

                if ocr_doc is None:
                    ocr_doc = fitz.open(stream=file_bytes, filetype="pdf")
                pix = ocr_doc[i - 1].get_pixmap(dpi=200)
                text = _ocr_image_bytes(pix.tobytes("png")).strip()
            except Exception:
                text = ""

        if text:
            pages.append(RawPage(page_number=i, text=text))

    return pages


def _extract_docx(file_bytes: bytes) -> list[RawPage]:
    import docx

    document = docx.Document(io.BytesIO(file_bytes))
    text = "\n\n".join(p.text for p in document.paragraphs if p.text.strip())
    return [RawPage(page_number=None, text=text)] if text.strip() else []


def _extract_pptx(file_bytes: bytes) -> list[RawPage]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[RawPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        text = "\n\n".join(parts)
        if text.strip():
            pages.append(RawPage(page_number=i, text=text))
    return pages


def _extract_tabular(file_bytes: bytes, ext: str) -> list[RawPage]:
    """
    CSV/XLSX: convert each sheet (or the single CSV table) into a compact
    textual representation, one "page" per sheet. Row-by-row text keeps
    retrieval useful for "what's the value of X" style questions without
    needing a separate structured-data query path.
    """
    import pandas as pd

    pages: list[RawPage] = []
    if ext == "csv":
        df = pd.read_csv(io.BytesIO(file_bytes))
        pages.append(RawPage(page_number=1, text=df.to_string(index=False)))
    else:
        sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
        for i, (name, df) in enumerate(sheets.items(), start=1):
            text = f"Sheet: {name}\n{df.to_string(index=False)}"
            pages.append(RawPage(page_number=i, text=text))
    return pages


def _extract_html(file_bytes: bytes) -> list[RawPage]:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    text = soup.get_text(separator="\n\n")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return [RawPage(page_number=None, text=text)] if text else []


def _extract_image(file_bytes: bytes) -> list[RawPage]:
    text = _ocr_image_bytes(file_bytes).strip()
    return [RawPage(page_number=1, text=text)] if text else []


def extract_pages(file_bytes: bytes, content_type: str, filename: str) -> list[RawPage]:
    """Turn raw upload bytes into a list of (page_number, text), dispatched by extension."""
    ext = _file_ext(filename)

    if ext == "pdf":
        return _extract_pdf(file_bytes)
    if ext == "docx":
        return _extract_docx(file_bytes)
    if ext == "pptx":
        return _extract_pptx(file_bytes)
    if ext in ("csv", "xlsx", "xls"):
        return _extract_tabular(file_bytes, ext)
    if ext in ("html", "htm"):
        return _extract_html(file_bytes)
    if ext in ("png", "jpg", "jpeg", "webp"):
        return _extract_image(file_bytes)

    # txt, md, and anything else falls back to plain-text decoding
    text = file_bytes.decode("utf-8", errors="ignore")
    return [RawPage(page_number=None, text=text)] if text.strip() else []


_HEADING_RE = re.compile(r"^(#{1,6}\s+.+|[A-Z][A-Za-z0-9 ,'&-]{2,60})$")


def _split_paragraphs(text: str) -> list[tuple[str, str | None]]:
    """
    Split page text into (paragraph, section_title) pairs. A short,
    title-cased or markdown-heading line is treated as a new section
    heading and attached to the paragraphs that follow it.
    """
    raw_paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    results = []
    current_section = None
    for para in raw_paragraphs:
        first_line = para.splitlines()[0].strip()
        if len(para.splitlines()) == 1 and _HEADING_RE.match(first_line) and len(first_line.split()) <= 12:
            current_section = first_line.lstrip("# ").strip()
            continue
        results.append((para, current_section))
    return results


def chunk_document(
    pages: list[RawPage],
    chunk_size_tokens: int = 350,
    overlap_tokens: int = 60,
) -> list[Chunk]:
    """
    Pack paragraphs into token-bounded windows with overlap, preserving
    page numbers and best-effort section titles.
    """
    # Flatten to a stream of (paragraph_text, page_number, section_title)
    stream: list[tuple[str, int | None, str | None]] = []
    for page in pages:
        for para, section in _split_paragraphs(page.text):
            stream.append((para, page.page_number, section))

    if not stream:
        return []

    chunks: list[Chunk] = []
    current_paras: list[tuple[str, int | None, str | None]] = []
    current_tokens = 0

    def flush():
        if not current_paras:
            return
        content = "\n\n".join(p[0] for p in current_paras)
        page_numbers = [p[1] for p in current_paras if p[1] is not None]
        page_number = page_numbers[0] if page_numbers else None
        section_title = next((p[2] for p in current_paras if p[2]), None)
        chunks.append(Chunk(
            content=content,
            token_count=count_tokens(content),
            page_number=page_number,
            section_title=section_title,
        ))

    i = 0
    while i < len(stream):
        para, page_no, section = stream[i]
        para_tokens = count_tokens(para)

        # A single paragraph larger than the chunk budget gets hard-split.
        if para_tokens > chunk_size_tokens:
            flush()
            current_paras, current_tokens = [], 0
            words = para.split()
            step = max(1, int(len(words) * chunk_size_tokens / max(para_tokens, 1)))
            for start in range(0, len(words), step):
                piece = " ".join(words[start:start + step])
                chunks.append(Chunk(
                    content=piece,
                    token_count=count_tokens(piece),
                    page_number=page_no,
                    section_title=section,
                ))
            i += 1
            continue

        if current_tokens + para_tokens > chunk_size_tokens and current_paras:
            flush()
            # Build overlap: carry the tail paragraphs whose combined
            # tokens are <= overlap_tokens into the next chunk.
            overlap_paras = []
            overlap_count = 0
            for p in reversed(current_paras):
                t = count_tokens(p[0])
                if overlap_count + t > overlap_tokens:
                    break
                overlap_paras.insert(0, p)
                overlap_count += t
            current_paras = overlap_paras
            current_tokens = overlap_count

        current_paras.append((para, page_no, section))
        current_tokens += para_tokens
        i += 1

    flush()
    return chunks
